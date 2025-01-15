from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# プレゼンテーションを作成
prs = Presentation()

# スライド 1: タイトルスライド
slide_1 = prs.slides.add_slide(prs.slide_layouts[5])
title_box = slide_1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "大谷翔平: 二刀流の軌跡と未来"
title.font.bold = True
title.font.size = Pt(44)

subtitle = title_frame.add_paragraph()
subtitle.text = "野球界の新たな伝説"
subtitle.font.size = Pt(32)

# 図形: タイトルとサブタイトルを囲む枠
shape = slide_1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# スライド 2: はじめに
slide_2 = prs.slides.add_slide(prs.slide_layouts[5])
content_box = slide_2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "大谷翔平は、現代野球における革新者であり、彼の活躍はスポーツ界全体に影響を与えています。"
content.font.size = Pt(24)

# 図形: 写真のコラージュを示す枠
shape = slide_2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# スライド 3: 初期の軌跡
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])
content_box = slide_3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "岩手県奥州市での少年時代から、彼の才能は際立っていました。"
content.font.size = Pt(24)

# ポイント
points = [
    "小学2年生で野球を始める",
    "高校3年生で160 km/hを記録"
]
for point in points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)

# 図形: 少年時代の大谷の写真を示す枠
shape = slide_3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# スライド 4: 日本プロ野球での活躍
slide_4 = prs.slides.add_slide(prs.slide_layouts[5])
content_box = slide_4.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "北海道日本ハムファイターズでの「二刀流」挑戦"
content.font.size = Pt(24)

# ポイント
points = [
    "2013年プロ初勝利、初本塁打",
    "2014年「2桁勝利・2桁本塁打」達成",
    "2016年投手と指名打者でベストナイン受賞"
]
for point in points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)

# 図形: 日本ハム時代の試合写真を示す枠
shape = slide_4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# スライド 5: メジャーリーグへの挑戦
slide_5 = prs.slides.add_slide(prs.slide_layouts[5])
content_box = slide_5.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "ロサンゼルス・エンゼルスでの新たな挑戦"
content.font.size = Pt(24)

# ポイント
points = [
    "2018年MLB史上初の10登板&20HR&10盗塁",
    "2021年シーズンMVP受賞"
]
for point in points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)

# 図形: エンゼルスでの活躍シーンを示す枠
shape = slide_5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# スライド 6: 歴史的な記録
slide_6 = prs.slides.add_slide(prs.slide_layouts[5])
content_box = slide_6.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "大谷翔平が打ち立てた数々の記録"
content.font.size = Pt(24)

# ポイント
points = [
    "2022年「2桁勝利・2桁本塁打」",
    "2023年WBCでのMVP受賞",
    "2024年「50本塁打、50盗塁」達成"
]
for point in points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)

# 図形: 記録達成時の写真を示す枠
shape = slide_6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# スライド 7: ロサンゼルス・ドジャースへの移籍
slide_7 = prs.slides.add_slide(prs.slide_layouts[5])
content_box = slide_7.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "新たなステージでの挑戦"
content.font.size = Pt(24)

# ポイント
points = [
    "2023年12月、史上最高額の契約",
    "2024年ワールドシリーズ制覇に貢献"
]
for point in points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)

# 図形: ドジャースのユニフォームを着た大谷を示す枠
shape = slide_7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# スライド 8: 影響力と未来
slide_8 = prs.slides.add_slide(prs.slide_layouts[5])
content_box = slide_8.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "大谷翔平の影響力は野球界を超えて広がっています。"
content.font.size = Pt(24)

# ポイント
points = [
    "タイム誌「世界で最も影響力のある100人」",
    "スポーツ選手長者番付で世界5位"
]
for point in points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)

# 図形: 世界中のファンと交流する大谷の写真を示す枠
shape = slide_8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# スライド 9: 結論
slide_9 = prs.slides.add_slide(prs.slide_layouts[5])
content_box = slide_9.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "大谷翔平は、野球界の未来を切り開く存在です。"
content.font.size = Pt(24)

# 図形: 大谷の未来を象徴するような希望に満ちた写真を示す枠
shape = slide_9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# スライド 10: 質疑応答
slide_10 = prs.slides.add_slide(prs.slide_layouts[5])
content_box = slide_10.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "ご質問をお待ちしております。"
content.font.size = Pt(24)

# 図形: 大谷の笑顔の写真を示す枠
shape = slide_10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(9), Inches(3))
shape.fill.background()
line = shape.line
line.color.rgb = RGBColor(0, 0, 0)

# プレゼンテーションを保存
prs.save('大谷翔平_プレゼンテーション.pptx')
